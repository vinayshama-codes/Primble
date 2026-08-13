# -*- coding: utf-8 -*-
"""Build Primble-Investor-Deck.pptx from the slide content below.

Every slide is a native PowerPoint shape with real, editable text.
Edit the DECK list and re-run: py build_deck.py
"""
import re
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── palette ──────────────────────────────────────────────────────────────────
INK        = RGBColor(0x10, 0x14, 0x19)
INK2       = RGBColor(0x38, 0x41, 0x4D)
MUTED      = RGBColor(0x6A, 0x74, 0x82)
RULE       = RGBColor(0xD2, 0xD7, 0xDE)
RULE_SOFT  = RGBColor(0xE4, 0xE8, 0xED)
SURFACE    = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE2   = RGBColor(0xF7, 0xF8, 0xFA)
ACCENT     = RGBColor(0xE6, 0x1B, 0x84)
ACCENT_INK = RGBColor(0xA8, 0x12, 0x5F)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
W_60       = RGBColor(0xA6, 0xAD, 0xB6)
W_75       = RGBColor(0xC4, 0xC9, 0xCF)
OK         = RGBColor(0x0E, 0x7C, 0x5A)
WARN       = RGBColor(0xA9, 0x65, 0x00)

SERIF = "Georgia"
SANS  = "Segoe UI"
MONO  = "Consolas"

# ── geometry (inches) ────────────────────────────────────────────────────────
SW, SH   = 13.333, 7.5
M        = 0.62
CW       = SW - 2 * M          # 12.093 content width
HEAD_Y   = 0.40
RULE1_Y  = 0.72
H2_Y     = 0.95
WHY_Y    = 6.42


def build_runs(para, text, base_size, base_color, base_font, bold_color=None,
               bold_weight=True):
    """Render **bold** markup as separate runs so emphasis survives editing."""
    bold_color = bold_color or base_color
    for chunk in re.split(r"(\*\*[^*]+\*\*)", text):
        if not chunk:
            continue
        strong = chunk.startswith("**") and chunk.endswith("**")
        r = para.add_run()
        r.text = chunk[2:-2] if strong else chunk
        r.font.name = base_font
        r.font.size = Pt(base_size)
        r.font.bold = bold_weight if strong else False
        r.font.color.rgb = bold_color if strong else base_color


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def line(slide, x, y, w, color=RULE_SOFT, pt=0.75):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Emu(1))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def box(slide, x, y, w, h, fill=SURFACE2, border=RULE):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.10)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return tf


def para(tf, first, space_before=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if space_before:
        p.space_before = Pt(space_before)
    p.line_spacing = 1.16
    return p


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE CONTENT - edit here, then re-run this script
# ─────────────────────────────────────────────────────────────────────────────
DECK = [
 dict(kind="cover", act="Primble  ·  ACORD Submission Platform",
   h="Upload the file. Get a **carrier-ready** submission.",
   sub="Primble reads a broker's policy documents and returns completed ACORD forms, "
       "a quality score, and the exact list of what is still missing - with the client "
       "already asked.",
   boxes=[("ACORD forms","17","fully supported"),
          ("Fields mapped","5,800+","across all forms"),
          ("Underwriting rules","80+","checked every time"),
          ("Manual entry cut","Hours","down to minutes")]),

 dict(act="Market", h="A broker's day is retyping PDFs.",
   t="Commercial submissions arrive as declarations pages, loss runs and schedules. "
     "A human reads them and retypes the contents into ACORD forms - **one box at a "
     "time, hundreds of boxes per submission.**",
   boxes=[("The labour","A single form can run to a thousand fields. A multi-line risk needs five or six forms at once."),
          ("The rejection","Carriers bounce incomplete submissions. Every bounce costs a re-quote cycle and a day of market time."),
          ("The liability","A wrong limit on a certificate is not a typo. It is an **E&O claim** against the agency.")],
   why="The pain repeats on **every single submission** - which is exactly what a metered subscription is built for."),

 dict(act="The Product", h="One workflow, five moves.",
   t="Primble is not a form filler with AI bolted on. Every stage checks the one before "
     "it, and **the client is brought in to close whatever the documents cannot.**",
   flow=[("01","Upload","Drop the whole document package in, however it arrives."),
         ("02","Understand","Turn the pages into a structured picture of the risk."),
         ("03","Review","Surface conflicts, gaps and blockers before anything is generated."),
         ("04","Generate","Produce the right ACORD forms, filled and scored."),
         ("05","Close","Ask the insured only what is genuinely still missing.")],
   why="Competitors sell step 4. **Steps 3 and 5 are where the recurring value and the switching cost live.**"),

 dict(act="Understanding the Submission", h="Take the whole package, however it arrives.",
   t="Multiple files, mixed formats - declarations pages, loss runs, schedules, prior "
     "applications. Uploaded together and treated as **one submission.**",
   boxes=[("Reads all of it","Long policy documents are read end to end, not skimmed or sampled. Nothing is quietly skipped."),
          ("Guards the door","Files are checked for what they really are and scanned before anything is processed."),
          ("Shows real progress","The producer watches actual work happen, file by file - not a spinner on a timer.")],
   why="Handling the messy real-world package is why Primble can charge per **submission** rather than per page."),

 dict(act="Understanding the Submission", h="It reads pages that trip everything else up.",
   t="Insurance declarations pages are dense two-column layouts. Standard document readers "
     "routinely pair the wrong label with the wrong value on them - **handing you a tax ID "
     "where the carrier name should be.**",
   boxes=[("Detects","Primble recognises the signature of a misread page instead of trusting the output."),
          ("Repairs","It rebuilds the page's real structure and recovers the correct pairings."),
          ("Never risks","A page that was read correctly is left completely untouched. If a repair does not help, it is discarded.")],
   why="Anyone can call a document-reading service. **Knowing where those services fail** is domain knowledge earned on real broker files."),

 dict(act="Understanding the Submission", h="When one upload holds two different businesses.",
   t="Brokers drag folders, and folders hold other people's policies. Primble spots it and "
     "**stops before treating strangers as one insured** - then hands the producer the "
     "controls.",
   boxes=[("Verify","A dedicated screen shows what appears to belong together and what does not, document by document."),
          ("Split or remove","Drop the documents that do not belong, or split the upload into separate submissions."),
          ("Or carry on","Park a split submission for later, add a missing document, or confirm it is fine and continue.")],
   why="A wrong insured is the most expensive error in the pipeline, and this is **the cheapest possible place to catch it.**"),

 dict(act="Submission Review", h="Submission Review: what we found.",
   t="Before a single form is generated, the producer gets one screen showing what was read "
     "and **where the documents disagree with each other.**",
   boxes=[("Documents processed","Every file, what Primble understood it to be, and what it contributed. Anything mistyped can be corrected in a click."),
          ("Data inconsistencies","Where the same detail appears with two different values, both are shown with the document each came from."),
          ("Suggested, not decided","Primble marks the value it believes is the best fit - but the producer picks, and their choice flows into every form at once.")],
   why="The producer keeps authority over every number. **Primble advises; the broker decides - and that is what makes the output defensible.**"),

 dict(act="Submission Review", h="Blockers and warnings, with a fix button on each.",
   t="The same screen lists what is wrong and what it costs. **A blocker caps the quality "
     "score at 60. A warning caps it at 85.** Both lift the moment they are genuinely fixed.",
   boxes=[("Grouped and ranked","Findings are clustered and ordered so the producer knows what to fix first, not handed a wall of errors."),
          ("Fixed in place","Each finding opens the exact inputs that clear it. Fill them and the score moves while you watch."),
          ("Or dismissed, on record","A finding can be waived with a reason. The waiver is recorded, never silently dropped.")],
   why="Fixing problems **before** the submission goes to market is the whole product. Everything downstream gets easier because of this screen."),

 dict(act="Generating the Forms", h="It tells you which forms this risk needs.",
   t="The producer does not pick from a menu of seventeen. Primble reads the exposures and "
     "**recommends the form set, in tiers, with a reason attached to every recommendation.**",
   boxes=[("Required","Core forms this submission cannot go without."),
          ("Recommended","Supporting forms justified by the exposures actually found."),
          ("Optional","Include if the exposure applies to this account."),
          ("Needs confirmation","Confirm relevance before generating.")],
   pts=["Every recommendation shows **why** it was made - and the reasoning is kept on the record.",
        "Nothing is forced: the producer can add any form Primble did not suggest, and select as many as the submission needs."],
   why="Form selection is where junior staff lose the most time and make the most mistakes. **This is the first visible win in any demo.**"),

 dict(act="Generating the Forms", h="How a form actually gets filled.",
   t="Filling an ACORD form is three jobs, not one. Primble does them in order and "
     "**checks its own work at every step.**",
   flow=[("STEP 1","Understand once","The package is read a single time into a structured view of the risk, shared by every form in the submission."),
         ("STEP 2","Place what is certain","Details that map cleanly onto a known ACORD box go straight there, the same way every time."),
         ("STEP 3","Resolve the rest","Everything ambiguous is worked out against the documents, then verified before it is allowed onto the form.")],
   why="Because the reading happens once and is shared, **the sixth form in a submission costs almost nothing to produce.**"),

 dict(act="Generating the Forms", h="No proof, no answer.",
   t="Every yes/no question and every checkbox has to be backed by something actually said "
     "in the uploaded documents. **An answer that cannot be evidenced is left blank, not "
     "guessed at.**",
   boxes=[("Across every form","Applies to disclosure questions on every ACORD form Primble supports, not a hand-picked few."),
          ("No borrowed proof","One sentence cannot be stretched to justify twenty different answers. Primble catches that and clears them."),
          ("Checked on the way out","Values are also checked for the right shape - a vehicle number can never land in a tax-ID box.")],
   why="Blank is a question for the client. Wrong is an E&O claim. **Primble is built to always fail toward blank.**"),

 dict(act="Generating the Forms", h="Every field says where it came from.",
   t="The producer reviews the real form with **each box colour-coded by how much it should "
     "be trusted** - so attention goes exactly where it is needed instead of re-reading "
     "everything.",
   boxes=[("Still required","Empty and needed. Fill before sending."),
          ("Verify","Filled by AI but not confirmed in the documents. Check this one."),
          ("Confirmed","Filled by AI and found in the uploaded documents."),
          ("From the client","Answered directly by the insured in the questionnaire.")],
   pts=["A running count of each category sits above the form, so the producer knows how much work is genuinely left.",
        "Every value is editable in place, and an edit flows back into the score and the other forms immediately."],
   why="This is the trust layer. **It turns 'the AI filled it in' into 'I know which four boxes to check' - which is what makes a broker sign off.**"),

 dict(act="Generating the Forms", h="A large fleet is one table, not hundreds of questions.",
   t="Vehicles, drivers, locations and loss runs collapse into **a single editable grid per "
     "schedule, with spreadsheet import** - instead of one question per empty box.",
   boxes=[("Import","Drop in a CSV or spreadsheet and the whole schedule populates at once."),
          ("Validate","Rows are checked as they land, and duplicates are flagged before they reach a form."),
          ("Look up","Vehicle identification numbers are decoded automatically into make, model and year."),
          ("Never lose a row","Rows beyond what the printed form holds are kept and flagged, never silently dropped.")],
   why="Fleet and property schedules are exactly the accounts a broker will not do by hand. **This is where large accounts become winnable.**"),

 dict(act="Generating the Forms", h="It understands coverage, not just boxes.",
   t="Commercial auto policies encode which vehicles a coverage protects as a single digit. "
     "Primble knows what those digits mean and **reasons about them the way an underwriter "
     "does.**",
   boxes=[("Stops false alarms","It knows when broad coverage already includes narrower coverage - so it stops asking for something the policy has."),
          ("Reads it in context","A coverage code is captured tied to the coverage line it belongs to, because the number alone means nothing."),
          ("Finds real gaps","It flags the genuine exposure the old approach missed: a fleet whose liability coverage does not actually reach the trucks.")],
   why="This is not form filling, it is **coverage analysis** - the wedge from a productivity tool into an underwriting product."),

 dict(act="Scoring", h="SQS: a credit score for a submission.",
   t="Primble grades every submission out of 100 on underwriting readiness - **before it "
     "goes to market.** The broker learns what the carrier would have said while there is "
     "still time to fix it.",
   boxes=[("Grade A","90+","submission ready"),
          ("Grade B","80+","minor gaps"),
          ("Grade C","70+","needs work"),
          ("Below","D / F","do not send")],
   why="SQS is the **proprietary metric.** It is the number a broker checks daily, the reason they come back, and the one thing a competitor cannot copy from a screenshot."),

 dict(act="Scoring", h="Six pillars, weighted - and hard-capped.",
   t="The score is a weighted judgement, not a box-counting percentage, and a blocking "
     "defect **caps it outright** so a broken submission can never present as ready.",
   boxes=[("Completeness","25","is it all there"), ("Consistency","25","does it agree"),
          ("Property","15","risk detail"), ("Loss history","15","claims record"),
          ("Umbrella","10","limits stack up"), ("Narrative","10","reads well")],
   pts=["**A blocker caps the score at 60. A warning caps it at 85.** Caps lift live, the moment the underlying issue is genuinely resolved.",
        "The submission is scored as a whole **and** form by form - the package catches problems no single form can see."],
   why="A score that cannot be gamed by filling boxes is a score **carriers could eventually underwrite against.**"),

 dict(act="Scoring", h="80+ rules that read the submission as one risk.",
   t="Umbrella limits against the policies beneath them. Payroll against employee count. "
     "Loss runs against a no-loss declaration. **Errors that are only visible when you look "
     "at two forms at once.**",
   boxes=[("Cross-form checks","Rules that compare one form against another, running automatically once the form set is known."),
          ("Field and document checks","Format, range and plausibility rules on the values themselves."),
          ("Three severities","**Blockers** stop a submission. **Warnings** cap the score. **Advisories** inform without penalty.")],
   why="Every rule is a piece of codified underwriting judgement. That library **compounds**, and it is the part a new entrant cannot ship in a quarter."),

 dict(act="Closing the Gaps", h="What the documents cannot answer, the client answers.",
   t="Primble turns what is still missing into a short questionnaire, sends it to the "
     "insured on a secure link, and **writes their answers straight back into the forms "
     "and the score.**",
   boxes=[("No login","A private link. The client answers on a phone, in a browser, with nothing to install."),
          ("Only what they can answer","Internal form plumbing never reaches the client. They see business questions, grouped and prioritised."),
          ("'I'm not sure'","An honest escape hatch. It is never stamped on a form and becomes the producer's follow-up list instead.")],
   why="This puts Primble in front of the **broker's own client** - a second surface, a second brand impression, and a natural expansion path."),

 dict(act="Closing the Gaps", h="And the producer never loses the thread.",
   t="Sent, opened, in progress, submitted - the status of every outstanding questionnaire "
     "is **visible at a glance,** and survives long after the working session is closed.",
   boxes=[("Reminders","One click re-sends a pending questionnaire, with the nudge logged."),
          ("Client help","A built-in assistant explains any insurance term to the client in plain English."),
          ("A permanent record","What the client answered is stored once, encrypted, and never overwritten."),
          ("Full timeline","Every event on the submission, in order, from upload to download.")],
   why="The record of what a client disclosed is **evidence** - and it is the first thing an agency reaches for when a claim is disputed."),

 dict(act="Closing the Gaps", h="Pick up any submission, any time.",
   t="Work does not have to finish in one sitting. **Every submission can be reopened at "
     "whatever stage it was left** - reviewed, part-selected, forms generated or fully "
     "scored.",
   boxes=[("Nothing is lost","Uploads, corrections, chosen values, waived findings and client answers all persist with the submission."),
          ("Any stage","Come back before selecting forms, after generating them, or mid-way through the review. It opens where you left it."),
          ("Still live","The score and the findings recalculate on reopening, so what you see is current rather than a stale snapshot.")],
   why="Real submissions take days and wait on other people. **A tool that assumes one sitting gets abandoned; this one fits how brokers actually work.**"),

 dict(act="Delivery", h="Download exactly what you need.",
   t="One form, the summary alone, or the entire package - and **the producer is told what "
     "is still outstanding before the file is handed over,** never after.",
   boxes=[("A single form","Any one form on its own, with or without the analysis summary attached."),
          ("The summary","The submission overview and score explanation as a standalone document for the underwriter."),
          ("The whole package","Every generated form, the summary and the signature, bundled in one download.")],
   pts=["A final check runs before every download and lists anything still unresolved, so nothing leaves the building unnoticed.",
        "A genuinely incomplete form can still be downloaded - clearly marked as a draft, with the producer's reason recorded."],
   why="The summary document is what the **underwriter reads first.** Owning it is how Primble gets seen on the carrier side of the market."),

 dict(act="Delivery", h="It fits the systems the agency already runs on.",
   t="Primble is built to sit inside an agency's existing stack rather than beside it. Any "
     "submission can be handed off as **complete structured data** - not a PDF someone has "
     "to retype at the other end.",
   boxes=[("Everything travels","Applicant, dates, coverage lines, every filled field and the full score breakdown move together as one record."),
          ("A stable contract","The handoff format is versioned, so an agency's own systems can consume it and evolve independently."),
          ("Any scope","One form or an entire submission goes across the same way."),
          ("Built for the majors","Handoff to the two dominant agency management systems is built; the direct connection is the next integration.")],
   why="This is the **enterprise unlock.** The hard part - a clean, complete data contract - already exists; what remains is a partnership, not a rebuild."),

 dict(act="Trust & Platform", h="Everything is on the record, permanently.",
   t="Primble keeps a durable history of every decision made on a submission - not just the "
     "final PDF. **If a dispute arrives years later, the answer is retrievable.**",
   boxes=[("What was advised","Every recommendation and finding shown to the producer, with the score at that moment."),
          ("What was decided","Every value chosen between conflicting documents, every finding fixed, waived or reopened - and the reason given."),
          ("Where each value came from","For each field: which document it came from, how it was filled, and how confident Primble was."),
          ("What the client said","The insured's own answers, stored once and never altered, plus who downloaded what and when.")],
   why="Insurance runs on defensibility. **This is the file an agency's E&O carrier asks for** - and it is the reason Primble becomes hard to remove once it is in the workflow."),

 dict(act="Trust & Platform", h="Sensitive data, treated as sensitive.",
   t="The platform handles personal information, financials, loss history and signatures. "
     "**It was built for that from the start, not retrofitted.**",
   boxes=[("Encrypted","Client data, questionnaire responses and signatures are encrypted individually, not just stored behind a lock."),
          ("Controlled","Sign-in through Google, strict per-submission access checks, and separate administrator permissions."),
          ("Protected","Usage limits, upload ceilings and malware scanning enforced before anything is processed."),
          ("Monitored","Full operational logging and error tracking, with clean shutdown that never drops work in progress.")],
   why="Security gaps are what kill insurance deals at procurement. **This one was closed in advance.**"),

 dict(act="Trust & Platform", h="We drove the cost of delivery down, measurably.",
   t="AI processing is the raw material of this business, and it has been managed like one - "
     "**measured, governed, and cut without giving up any quality.**",
   boxes=[("Processing per submission","-70%","same output"),
          ("Cost per submission","-24%","measured, one release"),
          ("Quality","Held","verified, not assumed"),
          ("Governed","Every AI step is tracked, and no change ships without accounting for its cost and quality effect.")],
   why="At a few dollars of revenue per submission, **processing cost is gross margin.** This team has already proven it can move that line."),

 dict(act="Business Model", h="Simple plans. Pay for what you run.",
   t="One subscription, priced by how many submissions an agency puts through each month. "
     "**Unlimited users on every plan** - the whole agency works in it, not one licensed seat.",
   boxes=[("Essentials","$59","50 a month"),
          ("Professional","$129","100 a month"),
          ("Business","$449","400 a month"),
          ("Enterprise","Custom","volume + integrations")],
   pts=["Annual billing saves around a quarter. Going over the monthly allowance is charged per submission, so nobody is ever blocked mid-deal.",
        "Usage is counted when a submission is **analysed** - so the value is paid for when it is delivered, not only if a file gets downloaded."],
   why="Volume-based pricing means revenue **grows with the customer's own book** and needs no renegotiation to do it."),

 dict(act="The Case", h="Why this is very hard to copy.",
   t="A general-purpose AI and a PDF library will get anyone to a convincing demo in a "
     "weekend. **Not one of the following comes out of the model.**",
   pts=["**Codified underwriting judgement.** 80+ rules and a six-pillar score, each argued out against real broker submissions. A newcomer does not know which rules matter until they have processed thousands of files.",
        "**A failure catalogue earned in production.** Every safeguard exists because a specific, subtle, expensive error was seen on a real document and traced to its cause. That list cannot be designed from scratch - only accumulated.",
        "**The only product that reaches the insured.** Primble sits between the broker and their client and holds the permanent record of what was disclosed. That relationship, and that record, is not something a rival can retrofit.",
        "**Switching means losing the history.** Once an agency's decisions, disclosures and audit trail live here, leaving costs them their own defence file.",
        "**The economics are already proven.** Cost per submission down measurably with quality held - and the discipline in place to keep it there as volume grows."],
   why="The moat is **accumulated judgement plus accumulated record**, and both get deeper with every submission processed."),

 dict(act="The Case", h="Built. Now it needs users.",
   t="The platform runs end to end and has been validated on real broker documents. "
     "**What is left is deployment, distribution and the first paying cohort.**",
   boxes=[("Shipped","All 17 forms generating. Review, scoring, client questionnaire, audit trail and billing - live, working, tested end to end."),
          ("Next","Hosted deployment, direct connections into the major agency systems, and scale-out for very large submissions."),
          ("The ask","Capital to get from validated product to a paying cohort: hosting, a security review, and a go-to-market motion into independent agencies.")],
   why="The technical risk sits **behind** this business. What remains is commercial risk - which is exactly the risk investors are equipped to price."),
]


# ─────────────────────────────────────────────────────────────────────────────
def render_boxes(slide, boxes, top, height):
    n = len(boxes)
    gap = 0.22 if n <= 4 else 0.16
    w = (CW - gap * (n - 1)) / n
    for i, b in enumerate(boxes):
        x = M + i * (w + gap)
        tf = box(slide, x, top, w, height)
        p = para(tf, True)
        r = p.add_run()
        r.text = b[0].upper()
        r.font.name, r.font.size = MONO, Pt(8)
        r.font.color.rgb = MUTED
        if len(b) == 3:                      # big-number box
            p2 = para(tf, False, space_before=6)
            r2 = p2.add_run(); r2.text = b[1]
            r2.font.name, r2.font.size, r2.font.bold = MONO, Pt(26 if len(b[1]) <= 6 else 18), True
            r2.font.color.rgb = INK
            p3 = para(tf, False, space_before=4)
            r3 = p3.add_run(); r3.text = b[2]
            r3.font.name, r3.font.size = SANS, Pt(9.5)
            r3.font.color.rgb = MUTED
        else:
            p2 = para(tf, False, space_before=7)
            build_runs(p2, b[1], 11, INK2, SANS, bold_color=ACCENT_INK)


def render_flow(slide, flow, top):
    n = len(flow)
    gap = 0.24
    w = (CW - gap * (n - 1)) / n
    for i, (num, title, desc) in enumerate(flow):
        x = M + i * (w + gap)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(top),
                                     Inches(w), Pt(2.2))
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background(); bar.shadow.inherit = False
        tf = textbox(slide, x, top + 0.16, w, 1.9)
        p = para(tf, True)
        r = p.add_run(); r.text = num
        r.font.name, r.font.size, r.font.bold = MONO, Pt(9), True
        r.font.color.rgb = ACCENT
        p2 = para(tf, False, space_before=4)
        r2 = p2.add_run(); r2.text = title
        r2.font.name, r2.font.size, r2.font.bold = SANS, Pt(14), True
        r2.font.color.rgb = INK
        p3 = para(tf, False, space_before=5)
        r3 = p3.add_run(); r3.text = desc
        r3.font.name, r3.font.size = SANS, Pt(10)
        r3.font.color.rgb = MUTED


def render_pts(slide, pts, top):
    tf = textbox(slide, M, top, CW, 2.0)
    for i, s in enumerate(pts):
        p = para(tf, i == 0, space_before=0 if i == 0 else 9)
        sq = p.add_run(); sq.text = "■   "
        sq.font.name, sq.font.size = SANS, Pt(9)
        sq.font.color.rgb = ACCENT
        build_runs(p, s, 11.5, INK2, SANS, bold_color=INK)


def main(out_path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
    blank = prs.slide_layouts[6]
    total = len(DECK)

    for idx, s in enumerate(DECK, start=1):
        slide = prs.slides.add_slide(blank)
        cover = s.get("kind") == "cover"

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = INK if cover else SURFACE

        # header strip
        tf = textbox(slide, M, HEAD_Y, CW * 0.4, 0.26)
        p = para(tf, True)
        r = p.add_run(); r.text = f"{idx:02d}"
        r.font.name, r.font.size, r.font.bold = MONO, Pt(9), True
        r.font.color.rgb = ACCENT
        r2 = p.add_run(); r2.text = f"   /   {total}"
        r2.font.name, r2.font.size = MONO, Pt(9)
        r2.font.color.rgb = W_60 if cover else MUTED

        tf = textbox(slide, M + CW * 0.42, HEAD_Y, CW * 0.58, 0.26)
        p = para(tf, True); p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = s["act"].upper()
        r.font.name, r.font.size = MONO, Pt(9)
        r.font.color.rgb = W_75 if cover else INK2

        line(slide, M, RULE1_Y, CW, RGBColor(0x33, 0x3A, 0x44) if cover else RULE_SOFT)

        if cover:
            tf = textbox(slide, M, 1.30, CW * 0.80, 2.4)
            p = para(tf, True); p.line_spacing = 1.06
            build_runs(p, s["h"], 40, WHITE, SERIF, bold_color=ACCENT, bold_weight=False)
            tf = textbox(slide, M, 3.95, CW * 0.62, 1.3)
            p = para(tf, True); p.line_spacing = 1.38
            build_runs(p, s["sub"], 13, W_75, SANS, bold_color=WHITE)
            n = len(s["boxes"]); gap = 0.22
            w = (CW - gap * (n - 1)) / n
            for i, b in enumerate(s["boxes"]):
                x = M + i * (w + gap)
                tfb = box(slide, x, 5.45, w, 1.30,
                          fill=RGBColor(0x1A, 0x1F, 0x26), border=RGBColor(0x33, 0x3A, 0x44))
                p = para(tfb, True)
                r = p.add_run(); r.text = b[0].upper()
                r.font.name, r.font.size = MONO, Pt(8); r.font.color.rgb = W_60
                p2 = para(tfb, False, space_before=8)
                r2 = p2.add_run(); r2.text = b[1]
                r2.font.name, r2.font.size, r2.font.bold = MONO, Pt(26), True
                r2.font.color.rgb = WHITE
                p3 = para(tfb, False, space_before=4)
                r3 = p3.add_run(); r3.text = b[2]
                r3.font.name, r3.font.size = SANS, Pt(9.5); r3.font.color.rgb = W_60
            continue

        # headline
        tf = textbox(slide, M, H2_Y, CW * 0.74, 1.15)
        p = para(tf, True); p.line_spacing = 1.10
        build_runs(p, s["h"], 27, INK, SERIF, bold_color=INK, bold_weight=False)

        # thesis
        tf = textbox(slide, M, 2.22, CW * 0.86, 1.0)
        p = para(tf, True); p.line_spacing = 1.38
        build_runs(p, s["t"], 12.5, INK2, SANS, bold_color=INK)

        body_top = 3.32
        if s.get("flow"):
            render_flow(slide, s["flow"], body_top + 0.10)
        elif s.get("boxes") and s.get("pts"):
            render_boxes(slide, s["boxes"], body_top - 0.05, 1.18)
            render_pts(slide, s["pts"], body_top + 1.32)
        elif s.get("boxes"):
            render_boxes(slide, s["boxes"], body_top, 2.55)
        elif s.get("pts"):
            render_pts(slide, s["pts"], body_top)

        # why-it-matters footer
        line(slide, M, WHY_Y, CW, RULE_SOFT)
        tf = textbox(slide, M, WHY_Y + 0.14, 1.55, 0.3)
        p = para(tf, True)
        r = p.add_run(); r.text = "WHY IT MATTERS"
        r.font.name, r.font.size, r.font.bold = MONO, Pt(8), True
        r.font.color.rgb = ACCENT
        tf = textbox(slide, M + 1.65, WHY_Y + 0.12, CW - 1.65, 0.62)
        p = para(tf, True); p.line_spacing = 1.30
        build_runs(p, s["why"], 11, INK2, SANS, bold_color=INK)

    prs.save(out_path)
    print(f"wrote {out_path}  ({total} slides)")


if __name__ == "__main__":
    main(sys.argv[1] if len(sys.argv) > 1 else "Primble-Investor-Deck.pptx")
